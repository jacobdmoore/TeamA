from pptx import Presentation
prsname=input("Enter Prestation name to modify: ")
prs=Presentation(prsname)
numslides=len(prs.slides) # produces number of slides in deck
f=open("filenumber.txt",'r')
lines=f.readlines()
lines=int(lines[0]) # lines is a list, call first element and turn into a int
numfile_in_textFile=lines-3 #3 extra txt files are associated with a folder in ubuntu even if empty
print('Found '+str(numfile_in_textFile)+' notes to be added to slides')
print('There are '+str(numslides)+' slides in this presentation')
f.close()
textslidenum=1
slidenum=0
while textslidenum<=numfile_in_textFile:
    textfile_name='textFiles/slide'+str(textslidenum)+".txt"
    f=open(textfile_name,'r')
    #f=open('textFiles/slide2.txt','r')
    lines1=f.read()
    #for line in lines1:
	    #print (line)
    f.close()
    textslidenum+=1
    current_slide=prs.slides[slidenum]
    notes_slide=current_slide.notes_slide
     #notes_slide=slide.notes_slide : acceses nots slide
    text_frame=notes_slide.notes_text_frame
   # print('Current Slide '+str(slidenum+1))
    newtext=lines1
    text_frame.text=text_frame.text+newtext
    slidenum+=1

newprs=prsname.split('.') #removes pptx fromname
new=newprs[0]
prs.save(new+' with notes.pptx') #adds withnotes.pptx to name


# print(first_slide.has_notes_slide)
# slide.has_notes_slide #checks to see if slide has notes slide