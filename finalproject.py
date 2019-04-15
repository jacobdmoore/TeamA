from pptx import Presentation
prsname=raw_input("Enter Prestation name to modify: ")
prs=Presentation(prsname)
numslides=len(prs.slides) # produces number of slides in deck
slidenum=0

while slidenum<numslides:

    current_slide=prs.slides[slidenum]
    notes_slide=current_slide.notes_slide
    #notes_slide=slide.notes_slide : acceses nots slide
    text_frame=notes_slide.notes_text_frame
    print('Current Slide'+str(slidenum+1))
    newtext=raw_input("Enter text for slide here")
    text_frame.text=text_frame.text+newtext
    slidenum=slidenum+1

newprs=prsname.split('.') #removes pptx fromname
new=newprs[0]
prs.save(new+' with notes.pptx') #adds withnotes.pptx to name


#print(first_slide.has_notes_slide)
#slide.has_notes_slide #checks to see if slide has notes slide