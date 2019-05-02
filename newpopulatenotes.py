from pptx import Presentation
import os
from os import listdir
from os.path import isfile, join

pptx_list = []
for file in os.listdir(os.getcwd() + '/'):
    if file.endswith('.pptx'):
        pptx_list.append(file)

if len(pptx_list) == 1:
    found_pptx = input('Is the .pptx file you are using '+pptx_list[0]+'? [y/N]')
    if found_pptx in ['y','Y','yes','YES','Yes']:
        found_pptx = 1
        pptx_file = pptx_list[0]
    else:
        pptx_file = input('Please type the name of the .pptx file you are using, including ".pptx".')

prs=Presentation(pptx_file)

text_file_list = [f for f in listdir(os.getcwd() + '/textFiles/') if isfile(join(os.getcwd() + '/textFiles/', f))]
slide_number_list = []

for file in text_file_list:

    slide_number = file[5::]
    slide_number = slide_number[:-4]
    slide_number_list.append(slide_number)

slidefile = zip(slide_number_list,text_file_list)

slide_number_list = [slide for slide,_ in sorted(slidefile)]

for number in slide_number_list:

    with open(os.getcwd()+'/textFiles/slide'+number+'.txt','r') as c_file:
        f = c_file.read()

    presentation_mode = 1

    try:
        number = int(number)
    except:
        number = int(number[:-1])
        presentation_mode = 0

    current_slide = prs.slides[int(number)]
    notes_slide = current_slide.notes_slide
     #notes_slide=slide.notes_slide : acceses nots slide
    text_frame = notes_slide.notes_text_frame
   # print('Current Slide '+str(slidenum+1))
    if not presentation_mode and f:
        f = '[Outside of presentation mode]: '+ f
    newtext = f
    text_frame.text=text_frame.text+newtext

newprs=pptx_file.split('.') #removes pptx fromname
new=newprs[0]
prs.save(new+'_notes.pptx') #adds withnotes.pptx to name