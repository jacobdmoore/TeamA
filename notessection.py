from pptx import Presentation
prs=Presentation('Midsemester Update.pptx')
#numslides=len(prs.slides) # produces number of slides in deck
first_slide=prs.slides[0]

print(first_slide.has_notes_slide)
#slide.has_notes_slide #checks to see if slide has notes slide

notes_slide=first_slide.notes_slide #creates note side for 1st slide

#notes_slide=slide.notes_slide : acceses nots slide
text_frame=notes_slide.notes_text_frame
text_frame.text="this is slide 1 in the notes section"
prs.save('test2.pptx')